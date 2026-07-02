import logging
import tempfile
from dataclasses import dataclass
from pathlib import Path

import cv2
import fitz  # PyMuPDF
import numpy as np
import torch
from pptx import Presentation
from pptx.util import Inches

from .config import PipelineConfig
from .inpainting import Inpainter
from .native_text import extract_native_text
from .ocr import OcrEngine
from .reconcile import reconcile_native_and_ocr
from .slide_builder import add_text_box
from .watermark import find_content_watermark_boxes, resolve_watermark_boxes, strip_watermark_texts

logger = logging.getLogger(__name__)


@dataclass
class ConversionResult:
    output_pptx_path: Path
    page_count: int


class PdfToPptxConverter:
    """Converts an image-flattened PDF into an editable PPTX: renders each page,
    extracts/OCRs its text, erases that text from the background via inpainting,
    and rebuilds it as native, positioned PowerPoint text boxes.

    Construction loads PaddleOCR and LaMa onto the GPU (or CPU, if unavailable),
    which takes real time and VRAM -- create one instance per process and reuse
    it across every document/page rather than constructing it per conversion.
    """

    def __init__(self, config: PipelineConfig = None):
        self.config = config or PipelineConfig()
        gpu = torch.cuda.is_available()
        logger.info(
            "CUDA available: %s%s",
            gpu,
            f" ({torch.cuda.get_device_name(0)})" if gpu else " -- falling back to CPU",
        )
        self.ocr = OcrEngine(lang=self.config.ocr_lang)
        self.inpainter = Inpainter()
        logger.info("LaMa inpainting device: %s", self.inpainter.device)

    def convert(self, pdf_path, output_pptx_path, background_dir=None, progress_callback=None) -> ConversionResult:
        """Run the full pipeline on one PDF.

        progress_callback, if given, is called as progress_callback(page_index_1_based,
        total_pages) after each page finishes -- useful for surfacing job progress in a
        long-running service.
        """
        pdf_path = Path(pdf_path)
        output_pptx_path = Path(output_pptx_path)
        if not pdf_path.exists():
            raise FileNotFoundError(f"Source PDF not found: {pdf_path}")

        with tempfile.TemporaryDirectory() as tmp_dir:
            bg_dir = Path(background_dir) if background_dir else Path(tmp_dir)
            bg_dir.mkdir(parents=True, exist_ok=True)

            doc = fitz.open(str(pdf_path))
            prs = Presentation()
            prs.slide_width = Inches(self.config.slide_w_in)
            prs.slide_height = Inches(self.config.slide_h_in)
            blank_layout = prs.slide_layouts[6]
            total_pages = len(doc)

            for page_idx in range(total_pages):
                page = doc[page_idx]
                pw, ph = page.rect.width, page.rect.height
                pix = page.get_pixmap(matrix=fitz.Matrix(self.config.zoom, self.config.zoom))
                img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
                if pix.n == 4:
                    img_rgb = cv2.cvtColor(img, cv2.COLOR_RGBA2RGB)
                elif pix.n == 1:
                    img_rgb = cv2.cvtColor(img, cv2.COLOR_GRAY2RGB)
                else:
                    img_rgb = img  # already RGB
                H, W = img_rgb.shape[:2]

                native_texts = extract_native_text(page, W, H, pw, ph, self.config)
                ocr_candidates = self.ocr.extract(img_rgb, W, H, self.config)
                native_texts, ocr_texts = reconcile_native_and_ocr(native_texts, ocr_candidates, self.config)

                watermark_boxes = resolve_watermark_boxes(self.config, W, H)
                content_watermark_boxes = find_content_watermark_boxes(
                    native_texts + ocr_texts, self.config.watermark_texts
                )
                native_texts = strip_watermark_texts(native_texts, watermark_boxes, self.config.watermark_texts)
                ocr_texts = strip_watermark_texts(ocr_texts, watermark_boxes, self.config.watermark_texts)
                all_texts = native_texts + ocr_texts

                clean_bg_rgb = self.inpainter.clean(
                    img_rgb, all_texts, W, H, self.config,
                    extra_boxes=watermark_boxes + content_watermark_boxes,
                )
                bg_path = bg_dir / f"page_{page_idx}.png"
                cv2.imwrite(str(bg_path), cv2.cvtColor(clean_bg_rgb, cv2.COLOR_RGB2BGR))

                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(str(bg_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
                for t in all_texts:
                    add_text_box(slide, t, W, H, self.config)

                logger.info(
                    "[%d/%d] native=%d ocr=%d", page_idx + 1, total_pages, len(native_texts), len(ocr_texts)
                )
                if progress_callback:
                    progress_callback(page_idx + 1, total_pages)

            output_pptx_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_pptx_path))
            logger.info("Saved: %s", output_pptx_path)

        return ConversionResult(output_pptx_path=output_pptx_path, page_count=total_pages)
